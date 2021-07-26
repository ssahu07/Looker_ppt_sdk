import csv
import sys
import textwrap
import time
import looker_sdk
import pptx
from looker_sdk import models
# import sdk_exceptions
from pptx import Presentation
from pptx.util import Inches

sdk = looker_sdk.init31(r'C:/Users/Sunil.sahu/Desktop/Python_Programs/Looker.ini')


def main():
    """Given a look title, find the corresponding look id and use
    it to render its image.
    $ python download_look.py "A good look" 1024 768 png
    """
    f = open(
        r'C:\Users\Sunil.sahu\PycharmProjects\PPT_Embed\Input.csv')  # Update csv path and connection same as PowerBI_Input.csv
    csv_f = csv.reader(f)
    for row in csv_f:
        for col in csv_f:
            look_title = col[1]
            image_width = int(sys.argv[2]) if len(sys.argv) > 2 else 545
            image_height = int(sys.argv[3]) if len(sys.argv) > 3 else 842
            image_format = sys.argv[4] if len(sys.argv) > 4 else "png"
            look = get_look(look_title)
            download_look(look, image_format, image_width, image_height)

    img_path = 'TK Rate Analyzer Adhoc Report LK.png'

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    pic = slide.shapes.add_picture(img_path, pptx.util.Inches(0.5), pptx.util.Inches(0.5),
                                   width=pptx.util.Inches(9), height=pptx.util.Inches(5))
    prs.save('Test.pptx')

def get_look(title: str) -> models.Look:
    title = title.lower()
    look = next(iter(sdk.search_looks(title=title)), None)
    # if not look:
    #     raise sdk_exceptions.NotFoundError(f"look '{title}' was not found")
    assert isinstance(look, models.Look)
    return look


def download_look(look: models.Look, result_format: str, width: int, height: int):
    """Download specified look as png/jpg"""
    assert look.id
    id = int(look.id)
    task = sdk.create_look_render_task(id, result_format, width, height,)

    # if not (task and task.id):
    #     raise sdk_exceptions.RenderTaskError(
    #         f"Could not create a render task for '{look.title}'"
    #     )

    # poll the render task until it completes
    elapsed = 0.0
    delay = 0.5  # wait .5 seconds
    while True:
        poll = sdk.render_task(task.id)
        if poll.status == "failure":
            print(poll)
            # raise sdk_exceptions.RenderTaskError(f"Render failed for '{look.title}'")
        elif poll.status == "success":
            break
        time.sleep(delay)
        elapsed += delay
    print(f"Render task completed in {elapsed} seconds")

    result = sdk.render_task_results(task.id)
    filename = f"{look.title}.{result_format}"
    with open(filename, "wb") as f:
        f.write(result)
    print(f"Look saved to '{filename}'")

main()
